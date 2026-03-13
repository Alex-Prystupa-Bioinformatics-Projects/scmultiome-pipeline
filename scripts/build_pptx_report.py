#!/usr/bin/env python3
# =============================================================================
# build_pptx_report.py  |  Auto-generate PowerPoint report from pipeline output
# =============================================================================
# Converts pipeline output PDFs into a ready-to-present .pptx deck.
# PDFs are converted to PIL Images in memory — no PNG files are written to disk.
#
# Usage:
#   uv run python scripts/build_pptx_report.py \
#       --report         clustering \
#       --project_prefix jm-multiome \
#       --samplesheet    configs/samplesheet.csv
#
# Output: output/reports/{prefix}-clustering-report.pptx
# =============================================================================

import argparse
import csv
import os
import re
import sys
from io import BytesIO

from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches, Pt

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_title_slide(prs, project_prefix):
    slide = blank_slide(prs)
    txb = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(1.5))
    tf = txb.text_frame
    tf.paragraphs[0].text = f"{project_prefix} — Clustering Report"
    tf.paragraphs[0].runs[0].font.size = Pt(36)
    tf.paragraphs[0].runs[0].font.bold = True


def add_section_slide(prs, title):
    slide = blank_slide(prs)
    txb = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(1.5))
    tf = txb.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].runs[0].font.size = Pt(32)
    tf.paragraphs[0].runs[0].font.bold = True


def add_plot_slide(prs, pdf_path, title=None):
    if not os.path.exists(pdf_path):
        print(f"  [WARN] Missing, skipping: {pdf_path}", file=sys.stderr)
        return

    images = convert_from_path(pdf_path, dpi=150)
    buf = BytesIO()
    images[0].save(buf, format="PNG")
    buf.seek(0)

    slide = blank_slide(prs)

    top_offset = Inches(0)
    if title:
        txb = slide.shapes.add_textbox(Inches(0.2), Inches(0.05), Inches(12.9), Inches(0.45))
        tf = txb.text_frame
        tf.paragraphs[0].text = title
        tf.paragraphs[0].runs[0].font.size = Pt(14)
        top_offset = Inches(0.5)

    # Scale image to fit available slide space, preserve aspect ratio
    img_w, img_h = images[0].size
    avail_w = float(SLIDE_W)
    avail_h = float(SLIDE_H) - float(top_offset)
    scale   = min(avail_w / img_w, avail_h / img_h)
    disp_w  = img_w * scale
    disp_h  = img_h * scale
    left    = (avail_w - disp_w) / 2
    top     = float(top_offset) + (avail_h - disp_h) / 2

    slide.shapes.add_picture(buf, left, top, width=disp_w, height=disp_h)


def detect_resolutions(prefix):
    d = os.path.join("output", "plots", f"{prefix}-umaps", "clustering")
    if not os.path.isdir(d):
        return []
    pattern = re.compile(rf"^{re.escape(prefix)}-umap-wnn-res([\d.]+)\.pdf$")
    res = [m.group(1) for f in os.listdir(d) if (m := pattern.match(f))]
    return sorted(res, key=float)


def read_metadata_cols(samplesheet):
    with open(samplesheet, newline="") as f:
        return [c for c in csv.DictReader(f).fieldnames if c not in ("SampleID", "path")]


def build_clustering_report(prefix, samplesheet, output_dir):
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    udir = os.path.join("output", "plots", f"{prefix}-umaps")
    adir = os.path.join("output", "plots", f"{prefix}-annotation")
    cdir = os.path.join(udir, "clustering")

    # 1. Title
    add_title_slide(prs, prefix)

    # 2. Dimensionality Reduction
    add_section_slide(prs, "Dimensionality Reduction")
    add_plot_slide(prs, os.path.join(udir, f"{prefix}-all-reduction-umaps.pdf"),
                   title="RNA | ATAC | WNN UMAP")

    # 3. QC on Embedding
    add_section_slide(prs, "QC on Embedding")
    add_plot_slide(prs, os.path.join(udir, "qc-metrics", f"{prefix}-umap-wnn-qc-metrics.pdf"),
                   title="QC metrics on WNN UMAP")

    # 4. Sample & Metadata
    add_section_slide(prs, "Sample & Metadata")
    for col in read_metadata_cols(samplesheet):
        add_plot_slide(prs, os.path.join(udir, "metadata", f"{prefix}-umap-wnn-{col}.pdf"),
                       title=col)

    # 5. Clustering
    add_section_slide(prs, "Clustering")
    add_plot_slide(prs, os.path.join(cdir, f"{prefix}-clustree.pdf"),
                   title="Clustree — cell redistribution across resolutions")

    for res in detect_resolutions(prefix):
        add_section_slide(prs, f"Resolution {res}")
        add_plot_slide(prs, os.path.join(cdir, f"{prefix}-umap-wnn-res{res}.pdf"),
                       title=f"WNN UMAP — clusters (res {res})")
        add_plot_slide(prs, os.path.join(adir, "dotplot", f"{prefix}-res{res}-dotplot.pdf"),
                       title=f"DotPlot — top markers per cluster (res {res})")
        add_plot_slide(prs, os.path.join(adir, "heatmap", f"{prefix}-res{res}-heatmap.pdf"),
                       title=f"Heatmap — top markers per cluster (res {res})")

    os.makedirs(output_dir, exist_ok=True)
    out_path = os.path.join(output_dir, f"{prefix}-clustering-report.pptx")
    prs.save(out_path)
    print(f"Report saved: {out_path}  ({len(prs.slides)} slides)")


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--report",         required=True)
    parser.add_argument("--project_prefix", required=True)
    parser.add_argument("--samplesheet",    default="configs/samplesheet.csv")
    parser.add_argument("--output_dir",     default="output/reports")
    args = parser.parse_args()

    if args.report == "clustering":
        build_clustering_report(args.project_prefix, args.samplesheet, args.output_dir)
    else:
        print(f"Unknown report type: {args.report}", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
